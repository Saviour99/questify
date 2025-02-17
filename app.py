from flask import Flask, render_template, request, make_response
from flask_sqlalchemy import SQLAlchemy
from flask_migrate import Migrate
import spacy
from collections import Counter
import random
import os
from pptx import Presentation
from PyPDF2 import PdfReader
from dotenv import load_dotenv
from datetime import datetime

# Load environment variables from .env file
load_dotenv()

# Debug: Print environment variables
print("DB_USERNAME:", os.getenv('DB_USERNAME'))
print("DB_PASSWORD:", os.getenv('DB_PASSWORD'))
print("DB_NAME:", os.getenv('DB_NAME'))
print("DB_HOST:", os.getenv('DB_HOST'))
print("DB_PORT:", os.getenv('DB_PORT'))

# Retrieve environment variables
DB_USERNAME = os.getenv('DB_USERNAME')
DB_PASSWORD = os.getenv('DB_PASSWORD')
DB_NAME = os.getenv('DB_NAME')
DB_HOST = os.getenv('DB_HOST', 'localhost') 
DB_PORT = os.getenv('DB_PORT', '3306') 


app = Flask(__name__)

# Configure the SQLAlchemy part of the app instance
app.config['SQLALCHEMY_DATABASE_URI'] = f'mysql+mysqlconnector://{DB_USERNAME}:{DB_PASSWORD}@{DB_HOST}:{DB_PORT}/{DB_NAME}'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False


print("Database URI:", app.config['SQLALCHEMY_DATABASE_URI'])

# Create the SQLAlchemy db instance
db = SQLAlchemy(app)

# Create the Flask-Migrate instance
migrate = Migrate(app, db)

# Define the Q&A model
class QA(db.Model):
    __tablename__ = 'qa_table'
    source_file_id = db.Column(db.Integer, primary_key=True, autoincrement=True)
    file_name = db.Column(db.String(100), nullable=False)
    question_id = db.Column(db.Integer, nullable=False)
    question = db.Column(db.String(9000), nullable=False)
    option_a = db.Column(db.String(400), nullable=False)
    option_b = db.Column(db.String(400), nullable=False)
    option_c = db.Column(db.String(400), nullable=False)
    option_d = db.Column(db.String(400), nullable=False)
    answer = db.Column(db.String(400), nullable=False)
    timestamp = db.Column(db.String(50), nullable=False)

# Load English tokenizer, tagger, parser, NER, and word vectors
nlp = spacy.load("en_core_web_sm")

# Generate MCQS with four options A-D and the correct answer
def generate_mcqs(text, num_questions=5, filename=""):
    if text is None:
        return []

    doc = nlp(text)
    sentences = [sent.text for sent in doc.sents]
    num_questions = min(num_questions, len(sentences))
    selected_sentences = random.sample(sentences, num_questions)
    mcqs = []

    for sentence in selected_sentences:
        sent_doc = nlp(sentence)
        nouns = [token.text for token in sent_doc if token.pos_ == "NOUN"]
        if len(nouns) < 2:
            continue
        noun_counts = Counter(nouns)
        if noun_counts:
            subject = noun_counts.most_common(1)[0][0]
            question_stem = sentence.replace(subject, "______")
            answer_choices = [subject]
            distractors = list(set(nouns) - {subject})
            while len(distractors) < 3:
                distractors.append("None of the above")
            random.shuffle(distractors)
            for distractor in distractors[:3]:
                answer_choices.append(distractor)
            random.shuffle(answer_choices)
            correct_answer = chr(64 + answer_choices.index(subject) + 1)

            # Storing the output in the database
            qa = QA(
                file_name=filename,
                question_id=len(mcqs) + 1,
                question=question_stem,
                option_a=answer_choices[0],
                option_b=answer_choices[1],
                option_c=answer_choices[2],
                option_d=answer_choices[3],
                answer=correct_answer,
                timestamp=datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            )
            db.session.add(qa)
            db.session.commit()
            mcqs.append((question_stem, answer_choices, correct_answer))

    return mcqs

@app.route("/")
@app.route("/home/")
def home():
    return render_template("index.html")

@app.route("/about/")
def about():
    return render_template("about.html")

@app.route("/contact/")
def contact():
    return render_template("contact.html")

@app.route("/admin/")
def admin():
    return render_template("admin.html")

@app.route("/mcqs/")
def mcqs():
    return render_template("qna.html")

@app.route('/upload/', methods=['GET', 'POST'])
def upload():
    # Get the file from the frontend and send the output to the frontend
    if request.method == 'POST':
        text = ""
        if 'file' in request.files:
            file = request.files['file']
            if file.filename:
                if file.filename.endswith('.pdf'):
                    text += extract_pdf(file)
                elif file.filename.endswith('.txt'):
                    text += file.read().decode('utf-8') #Reading the content of the text file
                elif file.filename.endswith('.pptx'):
                    text += extract_pptx(file)
            else:
                return "No file selected", 400
        else:
            text = request.form['text']

        if 'num_questions' in request.form:
            try:
                num_questions = int(request.form['num_questions'])
            except ValueError:
                return "Invalid number of questions", 400
        else:
            return "Number of questions is missing", 400

        mcqs = generate_mcqs(text, num_questions=num_questions, filename=file.filename)
        mcqs_with_index = [(i + 1, mcq) for i, mcq in enumerate(mcqs)]
        return render_template('qna.html', mcqs=mcqs_with_index)

    return render_template('upload.html')

# Extract the content from the PDF file
def extract_pdf(file):
    text = ""
    pdf_reader = PdfReader(file)
    for page_num in range(len(pdf_reader.pages)):
        page_text = pdf_reader.pages[page_num].extract_text()
        if page_text:
            text += page_text
    return text

# Extract the content from the PowerPoint Presentation
def extract_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

if __name__ == '__main__':
    app.run(debug=False, port=5000)